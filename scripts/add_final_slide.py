"""
Add methodology slide as slide 19 to Final_Presentation.pptx
"""

from pptx import Presentation

def add_final_slide():
    """Add methodology slide to end of Final_Presentation"""
    
    # Load both presentations
    main_prs = Presentation("presentation/Final_Presentation.pptx")
    method_prs = Presentation("presentation/Methodology_Slide.pptx")
    
    # Get the methodology slide
    method_slide = method_prs.slides[0]
    
    # Copy slide to main presentation
    # Get blank layout from main presentation
    blank_layout = main_prs.slide_layouts[6]
    
    # Add new slide
    new_slide = main_prs.slides.add_slide(blank_layout)
    
    # Copy all shapes from methodology slide
    for shape in method_slide.shapes:
        el = shape.element
        newel = type(el)(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    # Save
    main_prs.save("presentation/Final_Presentation.pptx")
    print(f"Added slide 19 to Final_Presentation.pptx")
    print(f"Total slides: {len(main_prs.slides)}")

if __name__ == "__main__":
    add_final_slide()

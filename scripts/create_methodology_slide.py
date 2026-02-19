"""
Create methodology slide explaining SD ratio calculation for individual ui
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_methodology_slide():
    """Create single-slide presentation with methodology explanation"""
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Add blank slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # UC Blue
    uc_blue = RGBColor(0, 48, 135)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = "VB posteriors for individual uᵢ:"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = uc_blue
    
    # First section: VB posteriors
    vb_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(9), Inches(0.7))
    vb_frame = vb_box.text_frame
    vb_frame.word_wrap = True
    
    p1 = vb_frame.paragraphs[0]
    p1.text = "results[[config]]$vb$mu_u"
    p1.font.name = "Consolas"
    p1.font.size = Pt(14)
    p1.level = 0
    
    p2 = vb_frame.add_paragraph()
    p2.text = " - means for each uᵢ"
    p2.font.size = Pt(14)
    p2.level = 0
    
    p3 = vb_frame.add_paragraph()
    p3.text = "results[[config]]$vb$Sigma_uu"
    p3.font.name = "Consolas"
    p3.font.size = Pt(14)
    p3.level = 0
    
    p4 = vb_frame.add_paragraph()
    p4.text = " - covariance matrix (diagonal = variances)"
    p4.font.size = Pt(14)
    p4.level = 0
    
    # Second section: Gibbs posteriors
    gibbs_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(9), Inches(1.6))
    gibbs_frame = gibbs_box.text_frame
    gibbs_frame.word_wrap = True
    
    g1 = gibbs_frame.paragraphs[0]
    g1.text = "Gibbs posteriors for individual uᵢ:"
    g1.font.bold = True
    g1.font.size = Pt(16)
    
    g2 = gibbs_frame.add_paragraph()
    g2.text = "results[[config]]$gibbs"
    g2.font.name = "Consolas"
    g2.font.size = Pt(14)
    g2.level = 0
    
    g3 = gibbs_frame.add_paragraph()
    g3.text = " - MCMC sample matrix"
    g3.font.size = Pt(14)
    g3.level = 0
    
    g4 = gibbs_frame.add_paragraph()
    g4.text = "For Q=5: 36000 samples × 10 columns (β₀, β₁, β₂, τₑ, τᵤ, u₁, u₂, u₃, u₄, u₅)"
    g4.font.size = Pt(13)
    g4.level = 0
    
    g5 = gibbs_frame.add_paragraph()
    g5.text = "For Q=10: 36000 samples × 15 columns"
    g5.font.size = Pt(13)
    g5.level = 0
    
    g6 = gibbs_frame.add_paragraph()
    g6.text = "For Q=20: 36000 samples × 25 columns"
    g6.font.size = Pt(13)
    g6.level = 0
    
    g7 = gibbs_frame.add_paragraph()
    g7.text = "For Q=50: 36000 samples × 55 columns"
    g7.font.size = Pt(13)
    g7.level = 0
    
    g8 = gibbs_frame.add_paragraph()
    g8.text = "For Q=100: 36000 samples × 105 columns"
    g8.font.size = Pt(13)
    g8.level = 0
    
    # Third section: Calculation
    calc_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.65), Inches(9), Inches(1.5))
    calc_frame = calc_box.text_frame
    calc_frame.word_wrap = True
    
    c1 = calc_frame.paragraphs[0]
    c1.text = "We can calculate:"
    c1.font.bold = True
    c1.font.size = Pt(16)
    
    c2 = calc_frame.add_paragraph()
    c2.text = "SD_VB(uᵢ) = √Sigma_uu[i,i]"
    c2.font.size = Pt(14)
    c2.level = 0
    
    c3 = calc_frame.add_paragraph()
    c3.text = "SD_Gibbs(uᵢ) = standard deviation of Gibbs samples for uᵢ"
    c3.font.size = Pt(14)
    c3.level = 0
    
    c4 = calc_frame.add_paragraph()
    c4.text = "rᵢ = SD_VB(uᵢ) / SD_Gibbs(uᵢ)  for each group"
    c4.font.size = Pt(14)
    c4.level = 0
    
    # Save
    output_path = "presentation/Methodology_Slide.pptx"
    prs.save(output_path)
    print(f"Created: {output_path}")
    
    return output_path

if __name__ == "__main__":
    create_methodology_slide()

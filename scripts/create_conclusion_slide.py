"""
Create Conclusion slide for VI presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add blank slide
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Add title
left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.6)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "Conclusion: Variational Bayes Trade-offs"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# Add content
left = Inches(0.5)
top = Inches(1.1)
width = Inches(9)
height = Inches(4)
content_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = content_box.text_frame
text_frame.word_wrap = True

# Key Findings
p = text_frame.paragraphs[0]
p.text = "Key Findings:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)
p.space_after = Pt(6)

findings = [
    "VB achieves 100–1200× speedup over Gibbs sampling",
    "Fixed effects (β) estimates: Excellent accuracy",
    "Variance components (τᵤ, σ²ᵤ): Systematic under-dispersion",
    "Under-dispersion worsens with fewer observations per group"
]

for finding in findings:
    p = text_frame.add_paragraph()
    p.text = f"  • {finding}"
    p.font.size = Pt(14)
    p.space_after = Pt(4)
    p.level = 0

# Practical Guidance
p = text_frame.add_paragraph()
p.text = "Practical Guidance:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)
p.space_after = Pt(6)
p.space_before = Pt(12)

guidance = [
    "Use VB for: Fixed effects, large-scale problems, initial exploration",
    "Use Gibbs/HMC for: Variance components, hierarchical uncertainty, final inference"
]

for item in guidance:
    p = text_frame.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(14)
    p.space_after = Pt(4)

# Active Research
p = text_frame.add_paragraph()
p.text = "Active Research:"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 102, 204)
p.space_after = Pt(6)
p.space_before = Pt(12)

research = [
    "Structured variational approximations (Blei et al., 2017)",
    "Expectation propagation methods",
    "Hybrid VB-MCMC approaches",
    "Adaptive variational families"
]

for item in research:
    p = text_frame.add_paragraph()
    p.text = f"  • {item}"
    p.font.size = Pt(14)
    p.space_after = Pt(4)

# Final statement
p = text_frame.add_paragraph()
p.text = "The under-dispersion issue is well-documented and remains an active area of research."
p.font.size = Pt(14)
p.font.italic = True
p.space_before = Pt(12)
p.alignment = PP_ALIGN.CENTER

# Save
output_path = "presentation/Conclusion_Slide.pptx"
prs.save(output_path)
print(f"✓ Conclusion slide created: {output_path}")

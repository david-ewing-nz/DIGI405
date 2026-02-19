#!/usr/bin/env python3
"""
Create Testing Methodology slide with M3 caveat
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

os.chdir(r"d:\github\VI1")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# UC Blue
UC_BLUE = RGBColor(0, 48, 135)
RED = RGBColor(220, 20, 60)

# Title
left = Inches(0.5)
top = Inches(0.4)
width = Inches(9)
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "Testing Methodology"
title_frame.paragraphs[0].font.size = Pt(40)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = UC_BLUE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Content
left = Inches(0.8)
top = Inches(1.6)
width = Inches(8.4)
height = Inches(3.5)
content_box = slide.shapes.add_textbox(left, top, width, height)
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Testing approach
p = content_frame.paragraphs[0]
p.text = "Comparison Strategy"
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = UC_BLUE
p.space_after = Pt(15)

p = content_frame.add_paragraph()
p.text = "• Model 1: Exact posterior (Normal-Inverse-Gamma) vs VB vs Gibbs"
p.font.size = Pt(18)
p.space_after = Pt(8)

p = content_frame.add_paragraph()
p.text = "• Model 2: Gibbs sampling (MCMC/NUTS) as gold standard vs VB"
p.font.size = Pt(18)
p.space_after = Pt(8)

p = content_frame.add_paragraph()
p.text = "• Varied group counts Q ∈ {5, 10, 20, 50, 100} to assess shrinkage effects"
p.font.size = Pt(18)
p.space_after = Pt(8)

p = content_frame.add_paragraph()
p.text = "• Measured SD ratios (VB/Gibbs) to quantify under-dispersion"
p.font.size = Pt(18)
p.space_after = Pt(25)

# M3 caveat
p = content_frame.add_paragraph()
p.text = "Important Note: Model 3 Implementation"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = RED
p.space_after = Pt(10)

p = content_frame.add_paragraph()
p.text = "Model 3 (hierarchical logistic) was implemented incorrectly and results are not reliable. Analysis focuses on Models 1 and 2."
p.font.size = Pt(18)
p.font.color.rgb = RED

# Save
prs.save("presentation/Testing_Methodology_Slide.pptx")
print("Testing Methodology slide created: presentation/Testing_Methodology_Slide.pptx")

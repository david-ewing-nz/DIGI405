"""Fix three slides that are going off the page - reduce font sizes"""

from pptx import Presentation
from pptx.util import Pt

def fix_vi_definition_slide():
    """Reduce font sizes in VI_Definition_Slide"""
    prs = Presentation("presentation/VI_Definition_Slide.pptx")
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Reduce body text
                    if run.font.size and run.font.size > Pt(20):
                        if run.font.size >= Pt(32):
                            continue  # Keep title
                        run.font.size = Pt(18)
    
    prs.save("presentation/VI_Definition_Slide.pptx")
    print("Fixed VI_Definition_Slide.pptx")


def fix_vi_challenges_slide():
    """Remove items 5 and 6, renumber, reduce font sizes"""
    prs = Presentation("presentation/VI_Challenges_Slide.pptx")
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # Remove and renumber paragraphs
            text_frame = shape.text_frame
            paragraphs_to_keep = []
            
            for para in text_frame.paragraphs:
                text = para.text.strip()
                
                # Skip items 5 and 6
                if text.startswith("5. Coordinate") or text.startswith("6. Hyperparameter"):
                    continue
                
                # Renumber item 7 to 5
                if text.startswith("7. Additional"):
                    para.text = text.replace("7. Additional", "5. Additional")
                
                paragraphs_to_keep.append(para)
            
            # Reduce font sizes
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size and run.font.size > Pt(20):
                        if run.font.size >= Pt(32):
                            continue  # Keep title
                        run.font.size = Pt(16)
    
    prs.save("presentation/VI_Challenges_Slide.pptx")
    print("Fixed VI_Challenges_Slide.pptx")


def fix_conclusion_slide():
    """Reduce font sizes in Conclusion_Slide"""
    prs = Presentation("presentation/Conclusion_Slide.pptx")
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Reduce all non-title text
                    if run.font.size and run.font.size > Pt(20):
                        if run.font.size >= Pt(32):
                            continue  # Keep titles
                        run.font.size = Pt(16)
    
    prs.save("presentation/Conclusion_Slide.pptx")
    print("Fixed Conclusion_Slide.pptx")


if __name__ == "__main__":
    fix_vi_definition_slide()
    fix_vi_challenges_slide()
    fix_conclusion_slide()
    print("All slides fixed")

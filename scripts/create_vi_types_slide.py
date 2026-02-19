#!/usr/bin/env python3
"""
Create VI Types (FFVB vs MFVB) slide
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

os.chdir(r"d:\github\VI1")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# UC Blue
UC_BLUE = RGBColor(0, 48, 135)
ORANGE = RGBColor(255, 140, 0)

# Title
left = Inches(0.5)
top = Inches(0.4)
width = Inches(9)
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "Two Approaches to Variational Inference"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = UC_BLUE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Content
left = Inches(0.8)
top = Inches(1.5)
width = Inches(8.4)
height = Inches(3.5)
content_box = slide.shapes.add_textbox(left, top, width, height)
content_frame = content_box.text_frame
content_frame.word_wrap = True

# FFVB
p = content_frame.paragraphs[0]
p.text = "Fixed-Form Variational Bayes (FFVB)"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = UC_BLUE
p.space_after = Pt(10)

p = content_frame.add_paragraph()
p.text = "Assumes a specific parametric form (e.g., Gaussian) for the entire posterior"
p.font.size = Pt(18)
p.space_after = Pt(20)

# MFVB
p = content_frame.add_paragraph()
p.text = "Mean-Field Variational Bayes (MFVB)"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = UC_BLUE
p.space_after = Pt(10)

p = content_frame.add_paragraph()
p.text = "Factorises the posterior into independent components"
p.font.size = Pt(18)
p.space_after = Pt(5)

p = content_frame.add_paragraph()
p.text = "Enables closed-form updates for conditionally conjugate models"
p.font.size = Pt(18)
p.space_after = Pt(25)

# Focus statement
p = content_frame.add_paragraph()
p.text = "This presentation focuses exclusively on Mean-Field VB"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ORANGE
p.alignment = PP_ALIGN.CENTER

# Save
prs.save("presentation/VI_Types_Slide.pptx")
print("VI Types slide created: presentation/VI_Types_Slide.pptx")

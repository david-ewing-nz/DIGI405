#!/usr/bin/env python3
"""
Create a PowerPoint presentation from slides and images.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
from pathlib import Path

# Set working directory
os.chdir(r"d:\github\VI1")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)  # 16:9 format

# Define slide data: (title, image_path, subtitle)
slides_data = [
    ("Bayes' Theorem", "presentation/Bayesian.png", "Prior, likelihood, posterior, and evidence"),
    ("Bayesian vs Variational Inference", "presentation/BI-VI.png", "Sampling vs optimisation"),
    ("Finding the Optimal q in Q-space", "presentation/q-space-visual-fixed.png", "Visualising the search for optimal distribution"),
    ("Coordinate Ascent VI", "presentation/coordinate_ascent.png", "Iterate factor updates to maximise ELBO"),
    ("Core Models Overview", "presentation/three_models_clean.png", "M1 linear, M2 hierarchical linear, M3 hierarchical logistic"),
    ("Exact Posterior Inference", "presentation/Slide_7B_Exact_Posterior_Inference.png", "Analytical solution for linear model"),
    ("Mean-Field Inference", "presentation/Slide7C_Mean_Field_Inference.png", "Coordinate ascent variational Bayes"),
    ("Under-dispersion in M2 Variance Component", "presentation/M2_tau_u_overlay_comparison.png", "VB vs Gibbs for τ_u: side-by-side comparison"),
    ("Configuration Comparison", "presentation/M2_tau_u_5panel_overlay.png", "VB vs Gibbs across all Q configurations"),
    ("PLACEHOLDER_SD_RATIOS", None, "SD Ratio Comparison"),  # Will insert PowerPoint table slide
    ("Computational Advantage: VB vs Gibbs", "presentation/timing_dashboard.png", "VB achieves 100-1200× speedup with model-specific scaling patterns"),
]

for slide_num, (title, img_path, subtitle) in enumerate(slides_data, 1):
    # Add blank slide
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Handle SD Ratios table slide specially
    if title == "PLACEHOLDER_SD_RATIOS":
        try:
            from pptx import Presentation as LoadPresentation
            sd_ratios_prs = LoadPresentation("presentation/SD_Ratios_Table_Slide.pptx")
            # Remove the blank slide we just added
            rId = prs.slides._sldIdLst[-1].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[-1]
            # Copy the SD ratios slide
            for sd_slide in sd_ratios_prs.slides:
                new_slide = prs.slides.add_slide(blank_slide_layout)
                for shape in sd_slide.shapes:
                    el = shape.element
                    new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
            print(f"Slide {slide_num}: Added SD Ratios PowerPoint table")
        except Exception as e:
            print(f"Slide {slide_num}: Error adding SD Ratios table: {e}")
        continue
    
    # Add image - check if it exists (full screen, no titles)
    img_full_path = img_path
    if os.path.exists(img_full_path):
        try:
            # For PDF files, skip (PowerPoint can't embed PDFs)
            if img_full_path.lower().endswith('.pdf'):
                print(f"Slide {slide_num}: Skipping PDF file {img_path} (PowerPoint limitation)")
                # Add note about PDF
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = f"Note: PDF file {img_path} should be inserted manually"
            else:
                # All images full screen
                slide.shapes.add_picture(img_full_path, Inches(0), Inches(0), width=Inches(10), height=Inches(5.625))
                print(f"Slide {slide_num}: Added {img_path}")
        except Exception as e:
            print(f"Slide {slide_num}: Error adding {img_path}: {e}")
    else:
        print(f"Slide {slide_num}: File not found: {img_full_path}")

# Add conclusion slide from separate file
try:
    from pptx import Presentation as LoadPresentation
    conclusion_prs = LoadPresentation("presentation/Conclusion_Slide.pptx")
    for slide in conclusion_prs.slides:
        # Copy slide
        slide_layout = prs.slide_layouts[6]
        new_slide = prs.slides.add_slide(slide_layout)
        for shape in slide.shapes:
            el = shape.element
            new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
    print("Added conclusion slide")
except Exception as e:
    print(f"Could not add conclusion slide: {e}")

# Add references slide from separate file
try:
    references_prs = LoadPresentation("presentation/References_Slide.pptx")
    for slide in references_prs.slides:
        # Copy slide
        slide_layout = prs.slide_layouts[6]
        new_slide = prs.slides.add_slide(slide_layout)
        for shape in slide.shapes:
            el = shape.element
            new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
    print("Added references slide")
except Exception as e:
    print(f"Could not add references slide: {e}")

# Save presentation
output_path = "presentation/Variational_Inference_Presentation.pptx"
prs.save(output_path)
print(f"\nPresentation created: {output_path}")

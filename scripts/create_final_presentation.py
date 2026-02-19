#!/usr/bin/env python3
"""
Create final presentation from individual slides
"""
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import os

os.chdir(r"d:\github\VI1")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Slide order
slides_content = [
    {"type": "image", "path": "presentation/Bayesian.png"},
    {"type": "image", "path": "presentation/BI-VI.png"},
    {"type": "pptx", "path": "presentation/Title_Slide.pptx"},
    {"type": "pptx", "path": "presentation/VI_Definition_Slide.pptx"},
    {"type": "pptx", "path": "presentation/VI_Challenges_Slide.pptx"},
    {"type": "pptx", "path": "presentation/VI_Types_Slide.pptx"},
    {"type": "image", "path": "presentation/Slide_7B_Exact_Posterior_Inference.png"},
    {"type": "image", "path": "presentation/Slide7C_Mean_Field_Inference.png"},
    {"type": "image", "path": "presentation/q-space-visual-fixed-twice..png"},
    {"type": "image", "path": "presentation/three_models_comparison.png"},
    {"type": "pptx", "path": "presentation/Testing_Methodology_Slide.pptx"},
    {"type": "image", "path": "presentation/M2_tau_u_overlay_comparison.png"},
    {"type": "image", "path": "presentation/M2_tau_u_5panel_overlay.png"},
    {"type": "pptx", "path": "presentation/Methodology_Slide.pptx"},
    {"type": "pptx", "path": "presentation/Cascade_Proof_Slides.pptx"},
    {"type": "image", "path": "presentation/timing_dashboard.png"},
    {"type": "pptx", "path": "presentation/SD_Ratios_Table_Slide.pptx"},
    {"type": "pptx", "path": "presentation/Conclusion_Slide.pptx"},
    {"type": "pptx", "path": "presentation/References_Slide.pptx"},
    {"type": "pptx", "path": "presentation/Thank_You_Slide.pptx"}
]

def add_image_slide(prs, image_path):
    """Add image slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Get image dimensions
    img = Image.open(image_path)
    img_width, img_height = img.size
    img_ratio = img_width / img_height
    
    # Slide dimensions
    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches
    slide_ratio = slide_width / slide_height
    
    # Calculate dimensions to fit slide
    if img_ratio > slide_ratio:
        # Image wider than slide
        width = Inches(slide_width)
        height = Inches(slide_width / img_ratio)
        left = Inches(0)
        top = Inches((slide_height - height.inches) / 2)
    else:
        # Image taller than slide
        height = Inches(slide_height)
        width = Inches(slide_height * img_ratio)
        top = Inches(0)
        left = Inches((slide_width - width.inches) / 2)
    
    slide.shapes.add_picture(image_path, left, top, width, height)

def add_pptx_slide(prs, pptx_path):
    """Add slide from another PowerPoint file"""
    from copy import deepcopy
    source_prs = Presentation(pptx_path)
    
    for source_slide in source_prs.slides:
        # Use blank layout
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Copy all shapes from source slide
        for shape in source_slide.shapes:
            # Create deep copy of shape element
            el = shape.element
            newel = deepcopy(el)
            slide.shapes._spTree.append(newel)

# Add all slides
slide_num = 0
for content in slides_content:
    slide_num += 1
    try:
        if content["type"] == "image":
            if os.path.exists(content["path"]):
                add_image_slide(prs, content["path"])
                print(f"Slide {slide_num}: Added {content['path']}")
            else:
                print(f"Slide {slide_num}: Warning - {content['path']} not found")
        elif content["type"] == "pptx":
            if os.path.exists(content["path"]):
                add_pptx_slide(prs, content["path"])
                print(f"Slide {slide_num}: Added {content['path']}")
            else:
                print(f"Slide {slide_num}: Warning - {content['path']} not found")
    except Exception as e:
        print(f"Slide {slide_num}: Error - {e}")

# Save
output_path = "presentation/Final_Presentation.pptx"
prs.save(output_path)
print(f"\nFinal presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")

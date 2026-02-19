"""
Generate PowerPoint presentation with model rating slides
Extended version including comprehensive model comparisons
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation in 16:9 format
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Define slide content
slides_content = [
    {"type": "image", "path": "presentation/timing_vb_vs_gibbs_bars.png", "title": ""},
    {"type": "image", "path": "presentation/BI-VI.png", "title": ""},
    {"type": "image", "path": "presentation/title_slide.png", "title": ""},
    {"type": "image", "path": "presentation/q-space-visual-fixed-twice..png", "title": ""},
    {"type": "image", "path": "presentation/timing_dashboard.png", "title": ""},
    {"type": "image", "path": "presentation/comparison_sd_ratios_heatmap.png", "title": ""},
    {"type": "image", "path": "presentation/model_rating_heatmap.png", "title": ""},
    {"type": "image", "path": "presentation/model_rating_comparison.png", "title": ""},
    {"type": "image", "path": "presentation/M1_posteriors_comparison.png", "title": ""},
    {"type": "image", "path": "presentation/M2_tau_u_overlay_comparison.png", "title": ""},
    {"type": "image", "path": "presentation/M3_tau_u_overlay_comparison.png", "title": ""},
    {"type": "bullets", "title": "Parameter Types: What Can We Measure?", "bullets": [
        ("Fixed Effects (β)", [
            "Direct impact of predictors on outcomes",
            "Predict individual outcomes and estimate effect sizes",
            "Answer: 'How much does X change Y?'",
            "VI Performance: Excellent to Good (0.85-0.95)"
        ]),
        ("Observation Variance (τ_e or σ²)", [
            "Residual variation within groups",
            "Assess model fit and prediction uncertainty",
            "Answer: 'How precise are our predictions?'",
            "VI Performance: Good to Moderate (0.75-0.85)"
        ]),
        ("Variance Components (τ_u or σ²_u)", [
            "Between-group heterogeneity in hierarchical models",
            "Assess if groups meaningfully differ",
            "Answer: 'Do groups actually differ from each other?'",
            "VI Performance: Poor to Very Poor (0.40-0.70)"
        ])
    ]},
    {"type": "bullets", "title": "Model Capabilities and Limitations", "bullets": [
        ("Model 1 (Linear Regression)", [
            "✓ Excellent for fixed effects (0.90-0.95)",
            "✓ Good for observation variance (0.80-0.85)",
            "✗ Cannot estimate variance components (no hierarchical structure)",
            "Best for: Simple predictions without grouping"
        ]),
        ("Model 2 (Hierarchical Linear)", [
            "✓ Good for fixed effects across all configurations",
            "✓ Handles grouped/clustered data",
            "⚠ Performance degrades with more groups (Q=5 better than Q=100)",
            "✗ Poor variance component estimation at large Q",
            "Trade-off: 30-50% too narrow intervals for τ_u"
        ]),
        ("Model 3 (Hierarchical Logistic)", [
            "✓ Good for fixed effects (binary outcomes)",
            "✓ Provides probability predictions",
            "✗ Very poor for variance components (0.40-0.60)",
            "Critical: Do not trust VI for random effects variance"
        ])
    ]},
    {"type": "bullets", "title": "Key Findings: The VI Reliability Hierarchy", "bullets": [
        ("Parameter-Selective Failure", [
            "Fixed effects (β): 85-95% of correct width → VI reliable",
            "Observation variance (τ_e): 75-85% of correct width → VI acceptable",
            "Variance components (τ_u): 40-70% of correct width → VI unreliable",
            "Pattern: VI fails selectively as you move up the hierarchy"
        ]),
        ("Sample Size Effects", [
            "M2_Q5 (5 groups): Moderate performance for τ_u (0.75)",
            "M2_Q20 (20 groups): Caution required for τ_u (0.67)",
            "M2_Q100 (100 groups): Poor performance for τ_u (0.52)",
            "Larger hierarchical structures amplify under-dispersion"
        ]),
        ("Speed vs Accuracy Trade-off", [
            "VI provides ~100-150× speedup over HMC",
            "Fixed effects remain reliable despite speedup",
            "Variance components sacrifice accuracy for speed",
            "Must choose: fast inference or accurate uncertainty?"
        ])
    ]},
    {"type": "bullets", "title": "Practical Recommendations", "bullets": [
        ("When to Use VI Confidently", [
            "Model 1 for any purpose",
            "Model 2/3 for fixed effects and predictions",
            "Exploratory analysis requiring fast iteration",
            "When computational resources are limited"
        ]),
        ("When to Use HMC Instead", [
            "Model 2/3 for variance component inference",
            "Critical decisions about between-group variation",
            "Questions like 'How much do groups differ?'",
            "Regulatory requirements or high-stakes inference"
        ]),
        ("Hybrid Approach", [
            "Use VI for initial exploration and fixed effects",
            "Verify variance components with HMC spot-checks",
            "Document which parameters used which method",
            "Practical rule: VI for means, HMC for variances"
        ])
    ]},
    {"type": "bullets", "title": "Real-World Example: Housing Prices", "bullets": [
        ("Question 1: Effect of Room Size", [
            "Uses fixed effect (β_rm)",
            "VI estimate: β = 3.2 [2.8, 3.6]",
            "HMC estimate: β = 3.2 [2.7, 3.7]",
            "Conclusion: VI reliable for predictions"
        ]),
        ("Question 2: Neighbourhood Variation", [
            "Uses variance component (τ_u)",
            "VI estimate: τ_u = 0.5 [0.4, 0.6] ← 40% too narrow",
            "HMC estimate: τ_u = 0.5 [0.3, 0.8] ← realistic",
            "Conclusion: VI unreliable for between-group inference"
        ]),
        ("Key Insight", [
            "Same model, different reliability by parameter",
            "Must match method to question being asked",
            "Speed advantage only valuable if answer is trustworthy"
        ])
    ]}
]

# Helper function to add title
def add_title(slide, title_text):
    if title_text:
        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(28)
        title.text_frame.paragraphs[0].font.bold = True

# Add slides
for idx, content in enumerate(slides_content):
    if content["type"] == "image":
        # Blank layout for images
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Add title if specified (not for slides 1-2)
        if content["title"] and idx >= 2:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = content["title"]
            p = tf.paragraphs[0]
            p.font.size = Pt(28)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # Image below title
            img_top = Inches(0.9)
            img_height = Inches(4.525)
        else:
            # Full slide image
            img_top = 0
            img_height = Inches(5.625)
        
        # Add image
        try:
            if not content["title"]:
                # Full-screen: stretch to fill
                slide.shapes.add_picture(
                    content["path"],
                    left=0,
                    top=0,
                    width=Inches(10),
                    height=Inches(5.625)
                )
            else:
                # With title: maintain aspect ratio, center
                slide.shapes.add_picture(
                    content["path"],
                    left=Inches(0.5),
                    top=img_top,
                    width=Inches(9),
                    height=img_height
                )
        except Exception as e:
            print(f"Warning: Could not add image {content['path']}: {e}")
    
    elif content["type"] == "bullets":
        # Title and content layout
        bullet_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_layout)
        
        # Add title
        title = slide.shapes.title
        title.text = content["title"]
        title.text_frame.paragraphs[0].font.size = Pt(24)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Add content
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        for section_title, bullet_points in content["bullets"]:
            # Section heading
            p = tf.add_paragraph()
            p.text = section_title
            p.level = 0
            p.font.size = Pt(18)
            p.font.bold = True
            p.space_after = Pt(6)
            
            # Bullet points
            for bullet in bullet_points:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 1
                p.font.size = Pt(14)
                p.space_after = Pt(3)

# Save presentation
output_path = "presentation/VI_Analysis_Presentation_Complete.pptx"
prs.save(output_path)
print(f"\nPresentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("\nSlide outline:")
for idx, content in enumerate(slides_content, 1):
    if content["type"] == "image":
        print(f"  Slide {idx}: Image - {content.get('title', 'Full-screen image')}")
    else:
        print(f"  Slide {idx}: Bullets - {content['title']}")

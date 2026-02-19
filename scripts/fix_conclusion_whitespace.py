"""Remove whitespace and fix Conclusion_Slide layout"""

from pptx import Presentation
from pptx.util import Pt, Inches

prs = Presentation("presentation/Conclusion_Slide.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        text_frame = shape.text_frame
        
        # Remove paragraphs with cut-off footer text
        paragraphs = list(text_frame.paragraphs)
        for para in paragraphs:
            text = para.text.strip()
            if text.startswith("The under-dispersion issue"):
                # Clear this paragraph
                for run in para.runs:
                    run.text = ""
                para.text = ""
        
        # Reduce spacing between paragraphs
        for para in text_frame.paragraphs:
            para.space_before = Pt(0)
            para.space_after = Pt(4)

prs.save("presentation/Conclusion_Slide.pptx")
print("Removed whitespace from Conclusion_Slide.pptx")

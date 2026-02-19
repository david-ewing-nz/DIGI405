#!/usr/bin/env python3
"""
Create Thank You slide for Dr John Holmes
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

os.chdir(r"d:\github\VI1")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# UC Blue
UC_BLUE = RGBColor(0, 48, 135)

# Main thank you
left = Inches(1)
top = Inches(1.8)
width = Inches(8)
height = Inches(2.5)
content_box = slide.shapes.add_textbox(left, top, width, height)
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Thank you
p = content_frame.paragraphs[0]
p.text = "Thank You"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = UC_BLUE
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(30)

# Supervisor acknowledgement
p = content_frame.add_paragraph()
p.text = "Dr John Holmes"
p.font.size = Pt(36)
p.font.color.rgb = UC_BLUE
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(15)

# Role
p = content_frame.add_paragraph()
p.text = "Supervisor"
p.font.size = Pt(24)
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(20)

# Institution
p = content_frame.add_paragraph()
p.text = "School of Mathematics and Statistics"
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(5)

p = content_frame.add_paragraph()
p.text = "University of Canterbury"
p.font.size = Pt(20)
p.alignment = PP_ALIGN.CENTER

# Save
prs.save("presentation/Thank_You_Slide.pptx")
print("Thank You slide created: presentation/Thank_You_Slide.pptx")

#!/usr/bin/env python3
"""
Create VI Definition slide
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

os.chdir(r"d:\github\VI1")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# UC Blue
UC_BLUE = RGBColor(0, 48, 135)

# Title
left = Inches(0.5)
top = Inches(0.4)
width = Inches(9)
height = Inches(1)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "What is Variational Inference?"
title_frame.paragraphs[0].font.size = Pt(40)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = UC_BLUE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Content
left = Inches(1)
top = Inches(1.8)
width = Inches(8)
height = Inches(3)
content_box = slide.shapes.add_textbox(left, top, width, height)
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Main definition
p = content_frame.paragraphs[0]
p.text = "Variational Inference (VI) is an optimisation technique for Bayesian inference."
p.font.size = Pt(24)
p.font.bold = True
p.space_after = Pt(20)

# Explanation
p = content_frame.add_paragraph()
p.text = "Rather than computing the true posterior distribution directly (which is often intractable), VI approximates it by:"
p.font.size = Pt(20)
p.space_after = Pt(15)

# Bullet 1
p = content_frame.add_paragraph()
p.text = "1. Choosing a tractable family of distributions Q"
p.font.size = Pt(18)
p.level = 1
p.space_after = Pt(10)

# Bullet 2
p = content_frame.add_paragraph()
p.text = "2. Finding the member of Q that best approximates the true posterior"
p.font.size = Pt(18)
p.level = 1
p.space_after = Pt(10)

# Bullet 3
p = content_frame.add_paragraph()
p.text = "3. Measuring 'best' using KL divergence"
p.font.size = Pt(18)
p.level = 1
p.space_after = Pt(15)

# Key insight
p = content_frame.add_paragraph()
p.text = "VI transforms a difficult integration problem into a tractable optimisation problem."
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = UC_BLUE

# Save
prs.save("presentation/VI_Definition_Slide.pptx")
print("VI Definition slide created: presentation/VI_Definition_Slide.pptx")

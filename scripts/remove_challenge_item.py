"""Remove items 5 and 6 from VI_Challenges_Slide and renumber"""

from pptx import Presentation

prs = Presentation("presentation/VI_Challenges_Slide.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        text_frame = shape.text_frame
        
        for para in text_frame.paragraphs:
            text = para.text.strip()
            
            # Renumber item 6 to 5
            if text.startswith("6. Hyperparameter sensitivity"):
                para.text = text.replace("6. Hyperparameter", "5. Hyperparameter")
            
            # Renumber item 7 to 6
            elif text.startswith("7. Additional approximation"):
                para.text = text.replace("7. Additional", "6. Additional")

prs.save("presentation/VI_Challenges_Slide.pptx")
print("Updated VI_Challenges_Slide.pptx - removed items 5 & 6 and renumbered")

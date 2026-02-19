#!/usr/bin/env python3
"""
Create VI Challenges slide listing known issues
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# UC Blue
UC_BLUE = RGBColor(0, 48, 135)

# Title
left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.6)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "Variational Inference: Known Challenges"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = UC_BLUE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Content
left = Inches(1.5)
top = Inches(1.2)
width = Inches(7)
height = Inches(3.8)
content_box = slide.shapes.add_textbox(left, top, width, height)
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Challenge 1
p = content_frame.paragraphs[0]
p.text = "1. Under-dispersion (variance underestimation in posterior approximations)"
p.font.size = Pt(18)
p.space_after = Pt(12)

# Challenge 2
p = content_frame.add_paragraph()
p.text = "2. Computational sensitivity (learning rate tuning, initialization dependence)"
p.font.size = Pt(18)
p.space_after = Pt(12)

# Challenge 3
p = content_frame.add_paragraph()
p.text = "3. Local optima (convergence to suboptimal solutions)"
p.font.size = Pt(18)
p.space_after = Pt(12)

# Challenge 4
p = content_frame.add_paragraph()
p.text = "4. Approximation quality (mean-field independence assumptions)"
p.font.size = Pt(18)
p.space_after = Pt(12)

# Challenge 5
p = content_frame.add_paragraph()
p.text = "5. Coordinate ascent convergence (iteration requirements, convergence criteria)"
p.font.size = Pt(18)
p.space_after = Pt(12)

# Challenge 6
p = content_frame.add_paragraph()
p.text = "6. Hyperparameter sensitivity (prior specification, penalty terms)"
p.font.size = Pt(18)
p.space_after = Pt(12)

# Challenge 7
p = content_frame.add_paragraph()
p.text = "7. Additional approximation limitations (distributional family restrictions, structural constraints)"
p.font.size = Pt(18)
p.space_after = Pt(18)

# Focus statement
p = content_frame.add_paragraph()
p.text = "This presentation focuses specifically on challenge 1: under-dispersion."
p.font.size = Pt(20)
p.font.italic = True
p.font.bold = True
p.font.color.rgb = UC_BLUE

# Save
prs.save("presentation/VI_Challenges_Slide.pptx")
print("VI Challenges slide created: presentation/VI_Challenges_Slide.pptx")

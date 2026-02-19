"""Rebuild References_Slide.pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

uc_blue = RGBColor(0, 48, 135)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "References"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = uc_blue

# References text
ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(4.2))
ref_frame = ref_box.text_frame
ref_frame.word_wrap = True

references = [
    "Blei, D. M., Kucukelbir, A., & McAuliffe, J. D. (2017). Variational inference: A review for statisticians. Journal of the American Statistical Association, 112(518), 859–877.",
    "",
    "Holmes, J. (2024). Fixed-form and mean-field variational Bayes approaches. Unpublished manuscript.",
    "",
    "Stan Development Team. (2024). Stan: A probabilistic programming language. Retrieved from http://mc-stan.org",
    "",
    "For more information on under-dispersion in hierarchical models, see course materials and Dr. Holmes' recent work on cascade effects in variance component estimation."
]

for i, ref in enumerate(references):
    if i == 0:
        p = ref_frame.paragraphs[0]
    else:
        p = ref_frame.add_paragraph()
    
    p.text = ref
    p.font.size = Pt(14)
    p.space_after = Pt(6)
    p.level = 0

prs.save("presentation/References_Slide.pptx")
print("References_Slide.pptx rebuilt")

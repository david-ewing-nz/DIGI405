"""
Create References slide for VI presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add blank slide
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Add title
left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.5)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "References"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# Add content in two columns
left_col_left = Inches(0.4)
right_col_left = Inches(5.2)
top = Inches(0.9)
col_width = Inches(4.5)
height = Inches(4.5)

# Left column
left_box = slide.shapes.add_textbox(left_col_left, top, col_width, height)
left_frame = left_box.text_frame
left_frame.word_wrap = True

references_left = [
    {
        "authors": "Blei, D. M., Kucukelbir, A., & McAuliffe, J. D.",
        "year": "2017",
        "title": "Variational Inference: A Review for Statisticians",
        "journal": "Journal of the American Statistical Association, 112(518), 859-877"
    },
    {
        "authors": "Jordan, M. I., Ghahramani, Z., Jaakkola, T. S., & Saul, L. K.",
        "year": "1999",
        "title": "An Introduction to Variational Methods for Graphical Models",
        "journal": "Machine Learning, 37(2), 183-233"
    },
    {
        "authors": "Turner, R. E., & Sahani, M.",
        "year": "2011",
        "title": "Two Problems with Variational Expectation Maximisation for Time-Series Models",
        "journal": "Bayesian Time Series Models, 109-130"
    },
    {
        "authors": "Giordano, R., Broderick, T., & Jordan, M. I.",
        "year": "2018",
        "title": "Covariances, Robustness, and Variational Bayes",
        "journal": "Journal of Machine Learning Research, 19(51), 1-49"
    }
]

for ref in references_left:
    p = left_frame.paragraphs[0] if ref == references_left[0] else left_frame.add_paragraph()
    p.text = f"{ref['authors']} ({ref['year']}). {ref['title']}. {ref['journal']}"
    p.font.size = Pt(10)
    p.space_after = Pt(10)
    p.line_spacing = 1.1

# Right column
right_box = slide.shapes.add_textbox(right_col_left, top, col_width, height)
right_frame = right_box.text_frame
right_frame.word_wrap = True

references_right = [
    {
        "authors": "Hoffman, M. D., & Gelman, A.",
        "year": "2014",
        "title": "The No-U-Turn Sampler: Adaptively Setting Path Lengths in Hamiltonian Monte Carlo",
        "journal": "Journal of Machine Learning Research, 15(1), 1593-1623"
    },
    {
        "authors": "Carpenter, B., et al.",
        "year": "2017",
        "title": "Stan: A Probabilistic Programming Language",
        "journal": "Journal of Statistical Software, 76(1), 1-32"
    },
    {
        "authors": "Gelman, A., & Rubin, D. B.",
        "year": "1992",
        "title": "Inference from Iterative Simulation Using Multiple Sequences",
        "journal": "Statistical Science, 7(4), 457-472"
    },
    {
        "authors": "Bishop, C. M.",
        "year": "2006",
        "title": "Pattern Recognition and Machine Learning",
        "journal": "Springer"
    }
]

for ref in references_right:
    p = right_frame.paragraphs[0] if ref == references_right[0] else right_frame.add_paragraph()
    p.text = f"{ref['authors']} ({ref['year']}). {ref['title']}. {ref['journal']}"
    p.font.size = Pt(10)
    p.space_after = Pt(10)
    p.line_spacing = 1.1

# Save
output_path = "presentation/References_Slide.pptx"
prs.save(output_path)
print(f"✓ References slide created: {output_path}")

"""
Create standalone title slide PowerPoint
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# UC Blue
uc_blue = RGBColor(0, 48, 135)

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(0.5))
tf = title_box.text_frame
tf.text = "Under-dispersion in Variational Bayesian Inference"
p = tf.paragraphs[0]
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = uc_blue
p.alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.font.size = Pt(26)
p.font.color.rgb = uc_blue
p.alignment = PP_ALIGN.CENTER

# Author
author_box = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(8), Inches(0.3))
tf = author_box.text_frame
tf.text = "David Ewing"
p = tf.paragraphs[0]
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = RGBColor(51, 51, 51)
p.alignment = PP_ALIGN.CENTER

# Degree
degree_box = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(8), Inches(0.25))
tf = degree_box.text_frame
tf.text = "Masters of Applied Data Science"
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.italic = True
p.font.color.rgb = RGBColor(85, 85, 85)
p.alignment = PP_ALIGN.CENTER

# Supervisor
supervisor_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.2))
tf = supervisor_box.text_frame
tf.text = "Supervisor: Dr John Holmes"
p = tf.paragraphs[0]
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(85, 85, 85)
p.alignment = PP_ALIGN.CENTER

# Institution
inst1_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.2))
tf = inst1_box.text_frame
tf.text = "School of Mathematics and Statistics"
p = tf.paragraphs[0]
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(102, 102, 102)
p.alignment = PP_ALIGN.CENTER

inst2_box = slide.shapes.add_textbox(Inches(1), Inches(4.75), Inches(8), Inches(0.2))
tf = inst2_box.text_frame
tf.text = "University of Canterbury"
p = tf.paragraphs[0]
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(102, 102, 102)
p.alignment = PP_ALIGN.CENTER

# Date
date_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.2))
tf = date_box.text_frame
tf.text = "30 January 2026"
p = tf.paragraphs[0]
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(136, 136, 136)
p.alignment = PP_ALIGN.CENTER

# Save
prs.save("presentation/Title_Slide.pptx")
print("Standalone title slide created: presentation/Title_Slide.pptx")

#!/usr/bin/env python3
"""
Create Slide 11: SD Ratios table with interpretations
Uses hardcoded SD ratio data from comparison analysis
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Set working directory
os.chdir(r"d:\github\VI1")

# Hardcoded SD ratio data from comparison results
table_data = [
    {"Model": "M1", "Q": "—", "beta_0": 1.000, "beta_1": 0.993, "beta_2": 0.988, "tau_e": 0.994, "tau_u": "NA", "sigma2_e": "NA", "sigma2_u": "NA"},
    {"Model": "M2_Q5", "Q": 5, "beta_0": 0.721, "beta_1": 0.994, "beta_2": 1.014, "tau_e": 0.995, "tau_u": 0.817, "sigma2_e": "NA", "sigma2_u": "NA"},
    {"Model": "M2_Q10", "Q": 10, "beta_0": 0.889, "beta_1": 0.998, "beta_2": 0.994, "tau_e": 0.970, "tau_u": 0.850, "sigma2_e": "NA", "sigma2_u": "NA"},
    {"Model": "M2_Q20", "Q": 20, "beta_0": 0.961, "beta_1": 0.993, "beta_2": 0.988, "tau_e": 0.962, "tau_u": 0.801, "sigma2_e": "NA", "sigma2_u": "NA"},
    {"Model": "M2_Q50", "Q": 50, "beta_0": 0.992, "beta_1": 0.998, "beta_2": 0.985, "tau_e": 0.913, "tau_u": 0.658, "sigma2_e": "NA", "sigma2_u": "NA"},
    {"Model": "M2_Q100", "Q": 100, "beta_0": 0.994, "beta_1": 1.000, "beta_2": 0.986, "tau_e": 0.803, "tau_u": 0.372, "sigma2_e": "NA", "sigma2_u": "NA"},
]

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add blank slide
blank_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_layout)

# Define color mapping (from heat map)
def get_color(value):
    """Return RGB color based on SD ratio value"""
    if value == "NA" or value is None:
        return RGBColor(211, 211, 211)  # Grey for NA
    
    val = float(value)
    if val >= 0.95:
        return RGBColor(255, 255, 0)  # Yellow - good
    elif val >= 0.85:
        return RGBColor(255, 165, 0)  # Orange - caution
    elif val >= 0.70:
        return RGBColor(255, 99, 71)  # Red-orange - poor
    else:
        return RGBColor(220, 20, 60)  # Dark red - very poor

# Define interpretations for ALL rows
interpretations = {
    "M1": {
        "overall": ("All β near 1.0 = excellent", "Fixed effects highly reliable for predictions", "Trust for all inference")
    },
    "M2_Q5": {
        "tau_u": ("β₀=0.721 (28% too narrow), τᵤ=0.817 (18% too narrow)", "Intercept unreliable; severely underestimates group differences", "Do not trust for between-group questions")
    },
    "M2_Q10": {
        "tau_u": ("β₀=0.889 (11% too narrow), τᵤ=0.850 (15% too narrow)", "Moderate under-dispersion in intercept and group variation", "Use caution for variance component inference")
    },
    "M2_Q20": {
        "tau_u": ("τᵤ=0.801 (20% too narrow)", "Group variation underestimated", "Important concern for hierarchical inference")
    },
    "M2_Q50": {
        "tau_e": ("τᵤ=0.658 (34% too narrow)", "Serious under-dispersion in group variance", "Major problem for between-group uncertainty")
    },
    "M2_Q100": {
        "tau_u": ("τᵤ=0.372 (63% too narrow)", "Catastrophic under-dispersion", "Never use VI for this parameter - use HMC")
    }
}

# Create table: 7 rows (header + 6 data) x 11 columns (original 9 + 2 new)
rows = len(table_data) + 1
cols = 11  # Model, Q, β₀, β₁, β₂, τ_e, τ_u, σ²_e, σ²_u, Numerical+Practical, When

# Position: leave small margins
left = Inches(0.2)
top = Inches(0.3)
width = Inches(9.6)
height = Inches(5)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths (11 columns total) - adjusted for better readability
col_widths = [Inches(1.3), Inches(0.6), Inches(0.6), Inches(0.6), Inches(0.6), 
              Inches(0.6), Inches(0.6), Inches(0.6), Inches(0.6),
              Inches(1.9), Inches(1.3)]
for i, width in enumerate(col_widths):
    table.columns[i].width = width

# Header row
headers = ["Model", "Q", "β₀", "β₁", "β₂", "τₑ", "τᵤ", "σ²ₑ", "σ²ᵤ", 
           "What it means / Practical consequence", "When it matters"]
for col, header in enumerate(headers):
    cell = table.cell(0, col)
    cell.text = header
    cell.text_frame.paragraphs[0].font.size = Pt(14)  # Larger for Greek chars
    cell.text_frame.paragraphs[0].font.bold = True
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Data rows
for row_idx, row_data in enumerate(table_data, start=1):
    model = row_data["Model"]
    
    # Column 0: Model
    cell = table.cell(row_idx, 0)
    cell.text = str(model)
    cell.text_frame.paragraphs[0].font.size = Pt(11)
    cell.text_frame.paragraphs[0].font.bold = True
    
    # Column 1: Q
    cell = table.cell(row_idx, 1)
    q_val = row_data["Q"]
    cell.text = str(q_val)
    cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    # Columns 2-8: SD ratios with color coding
    param_cols = ["beta_0", "beta_1", "beta_2", "tau_e", "tau_u", "sigma2_e", "sigma2_u"]
    for col_idx, param in enumerate(param_cols, start=2):
        cell = table.cell(row_idx, col_idx)
        value = row_data[param]
        
        if value == "NA":
            cell.text = "NA"
        else:
            cell.text = f"{float(value):.3f}"
        
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.fill.solid()
        cell.fill.fore_color.rgb = get_color(value)
    
    # Columns 9-10: Interpretations for ALL rows
    if model in interpretations:
        # Get interpretation for this model
        if "overall" in interpretations[model]:
            # M1 case - single overall interpretation
            numerical, practical, when = interpretations[model]["overall"]
        else:
            # M2 cases - get specific parameter interpretation
            param_key = list(interpretations[model].keys())[0]
            numerical, practical, when = interpretations[model][param_key]
        
        # Column 9: Combined numerical + practical
        cell = table.cell(row_idx, 9)
        cell.text = f"{numerical}: {practical}"
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.word_wrap = True
        
        # Column 10: When it matters
        cell = table.cell(row_idx, 10)
        cell.text = when
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.word_wrap = True

# Save presentation
output_path = "presentation/SD_Ratios_Table_Slide.pptx"
prs.save(output_path)
print(f"\nSD Ratios table created: {output_path}")
print(f"Table dimensions: {rows} rows x {cols} columns")

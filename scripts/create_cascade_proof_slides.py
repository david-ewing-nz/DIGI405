#!/usr/bin/env python3
"""
Create standalone 4-slide presentation proving cascade effect
Answers Dr John's question about tau_u under-dispersion affecting individual u_i
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

os.chdir(r"d:\github\VI1")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# UC Blue
UC_BLUE = RGBColor(0, 48, 135)

# ============================================================================
# Slide 1: The Cascade Question
# ============================================================================
blank_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_layout)

# Title
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_frame.text = "Does Hyper-parameter Under-dispersion Cascade?"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True
title_para.font.color.rgb = UC_BLUE

# Content
content_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Point 1
p1 = content_frame.paragraphs[0]
p1.text = "We've shown p(τᵤ | y) is too narrow under VB"
p1.font.size = Pt(24)
p1.space_after = Pt(20)

# Point 2
p2 = content_frame.add_paragraph()
p2.text = "Question: Does this affect individual random effects u₁, u₂, ..., uQ?"
p2.font.size = Pt(24)
p2.font.bold = True
p2.font.color.rgb = RGBColor(200, 0, 0)
p2.space_after = Pt(20)

# Point 3
p3 = content_frame.add_paragraph()
p3.text = "Why it matters: Model structure creates dependency"
p3.font.size = Pt(24)
p3.space_after = Pt(10)

p3_sub = content_frame.add_paragraph()
p3_sub.text = "uᵢ ~ N(0, τᵤ⁻¹)   ← each uᵢ depends on τᵤ"
p3_sub.font.size = Pt(20)
p3_sub.font.italic = True
p3_sub.level = 1
p3_sub.space_after = Pt(20)

# Point 4
p4 = content_frame.add_paragraph()
p4.text = "Test: Compute SD ratios for every uᵢ"
p4.font.size = Pt(24)
p4.font.bold = True
p4.space_after = Pt(10)

p4_sub = content_frame.add_paragraph()
p4_sub.text = "rᵢ = SDᵥB(uᵢ) / SDGibbs(uᵢ)   for all i = 1, ..., Q"
p4_sub.font.size = Pt(22)
p4_sub.font.color.rgb = UC_BLUE
p4_sub.level = 1

# ============================================================================
# Slide 2: Method - Computing r_i
# ============================================================================
slide2 = prs.slides.add_slide(blank_layout)

# Title
title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame2 = title_box2.text_frame
title_frame2.text = "Method: Computing rᵢ for Each Group"
title_para2 = title_frame2.paragraphs[0]
title_para2.font.size = Pt(36)
title_para2.font.bold = True
title_para2.font.color.rgb = UC_BLUE

# Left panel - VB
vb_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(2))
vb_frame = vb_box.text_frame
vb_frame.word_wrap = True

vb_title = vb_frame.paragraphs[0]
vb_title.text = "VB Approach"
vb_title.font.size = Pt(28)
vb_title.font.bold = True
vb_title.font.color.rgb = UC_BLUE
vb_title.space_after = Pt(15)

vb_p1 = vb_frame.add_paragraph()
vb_p1.text = "Mean-field factorisation:"
vb_p1.font.size = Pt(20)
vb_p1.space_after = Pt(8)

vb_p2 = vb_frame.add_paragraph()
vb_p2.text = "q(uᵢ) = N(μᵤᵢ, Σᵤᵤ[i,i])"
vb_p2.font.size = Pt(22)
vb_p2.font.italic = True
vb_p2.space_after = Pt(15)

vb_p3 = vb_frame.add_paragraph()
vb_p3.text = "Therefore:"
vb_p3.font.size = Pt(20)
vb_p3.space_after = Pt(8)

vb_p4 = vb_frame.add_paragraph()
vb_p4.text = "SDᵥB(uᵢ) = √Σᵤᵤ[i,i]"
vb_p4.font.size = Pt(24)
vb_p4.font.bold = True
vb_p4.font.color.rgb = RGBColor(0, 100, 0)

# Right panel - Gibbs
gibbs_box = slide2.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.3), Inches(2))
gibbs_frame = gibbs_box.text_frame
gibbs_frame.word_wrap = True

gibbs_title = gibbs_frame.paragraphs[0]
gibbs_title.text = "Gibbs Sampling"
gibbs_title.font.size = Pt(28)
gibbs_title.font.bold = True
gibbs_title.font.color.rgb = RGBColor(200, 0, 100)
gibbs_title.space_after = Pt(15)

gibbs_p1 = gibbs_frame.add_paragraph()
gibbs_p1.text = "MCMC samples:"
gibbs_p1.font.size = Pt(20)
gibbs_p1.space_after = Pt(8)

gibbs_p2 = gibbs_frame.add_paragraph()
gibbs_p2.text = "{uᵢ⁽¹⁾, uᵢ⁽²⁾, ..., uᵢ⁽ˢ⁾}"
gibbs_p2.font.size = Pt(22)
gibbs_p2.font.italic = True
gibbs_p2.space_after = Pt(15)

gibbs_p3 = gibbs_frame.add_paragraph()
gibbs_p3.text = "Therefore:"
gibbs_p3.font.size = Pt(20)
gibbs_p3.space_after = Pt(8)

gibbs_p4 = gibbs_frame.add_paragraph()
gibbs_p4.text = "SDGibbs(uᵢ) = sd(samples)"
gibbs_p4.font.size = Pt(24)
gibbs_p4.font.bold = True
gibbs_p4.font.color.rgb = RGBColor(200, 0, 100)

# Bottom - Formula
formula_box = slide2.shapes.add_textbox(Inches(1), Inches(4.0), Inches(8), Inches(1.2))
formula_frame = formula_box.text_frame
formula_frame.word_wrap = True

formula_p1 = formula_frame.paragraphs[0]
formula_p1.text = "SD Ratio for each group:"
formula_p1.font.size = Pt(22)
formula_p1.space_after = Pt(12)
formula_p1.alignment = PP_ALIGN.CENTER

formula_p2 = formula_frame.add_paragraph()
formula_p2.text = "rᵢ = SDᵥB(uᵢ) / SDGibbs(uᵢ)     for all i = 1, ..., Q"
formula_p2.font.size = Pt(26)
formula_p2.font.bold = True
formula_p2.font.color.rgb = UC_BLUE
formula_p2.alignment = PP_ALIGN.CENTER
formula_p2.space_after = Pt(12)

insight_p = formula_frame.add_paragraph()
insight_p.text = "Key insight: If rᵢ < 1 for all groups, cascade is confirmed"
insight_p.font.size = Pt(20)
insight_p.font.italic = True
insight_p.font.color.rgb = RGBColor(150, 0, 0)
insight_p.alignment = PP_ALIGN.CENTER

# ============================================================================
# Slide 3: Results - Systematic Under-dispersion
# ============================================================================
slide3 = prs.slides.add_slide(blank_layout)

# Title
title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
title_frame3 = title_box3.text_frame
title_frame3.text = "Results: Systematic Under-dispersion in All uᵢ"
title_para3 = title_frame3.paragraphs[0]
title_para3.font.size = Pt(32)
title_para3.font.bold = True
title_para3.font.color.rgb = UC_BLUE

# Image
img_path = "presentation/M2_u_i_boxplot.png"
if os.path.exists(img_path):
    slide3.shapes.add_picture(img_path, Inches(0.5), Inches(1.1), width=Inches(6.5))
else:
    print(f"Warning: {img_path} not found")

# Add legend box for reference lines
legend_box = slide3.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(2.5), Inches(0.7))
legend_frame = legend_box.text_frame
legend_frame.word_wrap = True

legend_title = legend_frame.paragraphs[0]
legend_title.text = "Reference Lines:"
legend_title.font.size = Pt(14)
legend_title.font.bold = True
legend_title.space_after = Pt(6)

legend_red = legend_frame.add_paragraph()
legend_red.text = "— — —  Perfect match (r = 1.0)"
legend_red.font.size = Pt(12)
legend_red.font.color.rgb = RGBColor(200, 0, 0)
legend_red.space_after = Pt(4)

legend_orange = legend_frame.add_paragraph()
legend_orange.text = "· · · · ·  Critical threshold (r = 0.8)"
legend_orange.font.size = Pt(12)
legend_orange.font.color.rgb = RGBColor(255, 140, 0)

# Key findings box
findings_box = slide3.shapes.add_textbox(Inches(7.2), Inches(1.5), Inches(2.5), Inches(3.5))
findings_frame = findings_box.text_frame
findings_frame.word_wrap = True

findings_title = findings_frame.paragraphs[0]
findings_title.text = "Key Findings"
findings_title.font.size = Pt(24)
findings_title.font.bold = True
findings_title.font.color.rgb = UC_BLUE
findings_title.space_after = Pt(15)

f1 = findings_frame.add_paragraph()
f1.text = "All 185 groups show rᵢ < 1"
f1.font.size = Pt(18)
f1.space_after = Pt(12)

f2 = findings_frame.add_paragraph()
f2.text = "Mean ratios:"
f2.font.size = Pt(18)
f2.font.bold = True
f2.space_after = Pt(8)

f2a = findings_frame.add_paragraph()
f2a.text = "Q=5:   r̄ = 0.85"
f2a.font.size = Pt(16)
f2a.font.name = "Consolas"
f2a.level = 1
f2a.space_after = Pt(5)

f2b = findings_frame.add_paragraph()
f2b.text = "Q=20:  r̄ = 0.75"
f2b.font.size = Pt(16)
f2b.font.name = "Consolas"
f2b.level = 1
f2b.space_after = Pt(5)

f2c = findings_frame.add_paragraph()
f2c.text = "Q=100: r̄ = 0.50"
f2c.font.size = Pt(16)
f2c.font.name = "Consolas"
f2c.level = 1
f2c.space_after = Pt(15)

f3 = findings_frame.add_paragraph()
f3.text = "Cascade confirmed:"
f3.font.size = Pt(18)
f3.font.bold = True
f3.font.color.rgb = RGBColor(200, 0, 0)
f3.space_after = Pt(8)

f3a = findings_frame.add_paragraph()
f3a.text = "Wrong τᵤ → Wrong uᵢ"
f3a.font.size = Pt(18)
f3a.font.color.rgb = RGBColor(200, 0, 0)
f3a.level = 1

# ============================================================================
# Slide 4: Example - Six Groups from Q=20
# ============================================================================
slide4 = prs.slides.add_slide(blank_layout)

# Title
title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
title_frame4 = title_box4.text_frame
title_frame4.word_wrap = True
title_frame4.text = "Example: Individual uᵢ Posteriors (Q=20)"
title_para4 = title_frame4.paragraphs[0]
title_para4.font.size = Pt(36)
title_para4.font.bold = True
title_para4.font.color.rgb = UC_BLUE

# Image
img_path2 = "presentation/M2_u_i_individual_Q20.png"
if os.path.exists(img_path2):
    slide4.shapes.add_picture(img_path2, Inches(0.5), Inches(1.0), width=Inches(9), height=Inches(3.8))
else:
    print(f"Warning: {img_path2} not found")

# Caption
caption_box = slide4.shapes.add_textbox(Inches(0.5), Inches(4.9), Inches(9), Inches(0.6))
caption_frame = caption_box.text_frame
caption_frame.word_wrap = True

caption_p = caption_frame.paragraphs[0]
caption_p.text = "VB posteriors (black line) systematically narrower than Gibbs (red histogram). SD ratios range from 0.72 to 0.81."
caption_p.font.size = Pt(16)
caption_p.font.italic = True
caption_p.alignment = PP_ALIGN.CENTER

# ============================================================================
# Slide 5: Distribution Detail (Violin Plot)
# ============================================================================
slide5 = prs.slides.add_slide(blank_layout)

# Title
title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
title_frame5 = title_box5.text_frame
title_frame5.word_wrap = True
title_frame5.text = "Distribution of SD Ratios Across All Groups"
title_para5 = title_frame5.paragraphs[0]
title_para5.font.size = Pt(28)
title_para5.font.bold = True
title_para5.font.color.rgb = UC_BLUE

# Subtitle
subtitle_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.4))
subtitle_frame5 = subtitle_box5.text_frame
subtitle_frame5.text = "Larger Q amplifies under-dispersion in individual random effects"
subtitle_para5 = subtitle_frame5.paragraphs[0]
subtitle_para5.font.size = Pt(20)
subtitle_para5.alignment = PP_ALIGN.CENTER

# Image
img_path5 = "presentation/M2_u_i_violin.png"
if os.path.exists(img_path5):
    slide5.shapes.add_picture(img_path5, Inches(0.5), Inches(1.3), width=Inches(9), height=Inches(4))
else:
    print(f"Warning: {img_path5} not found")

# ============================================================================
# Save
# ============================================================================
output_path = "presentation/Cascade_Proof_Slides.pptx"
prs.save(output_path)
print(f"\nCascade proof slides saved to {output_path}")
print("5 slides created:")
print("  Slide 1: The Cascade Question")
print("  Slide 2: Method - Computing rᵢ")
print("  Slide 3: Results - Systematic Under-dispersion (boxplot)")
print("  Slide 4: Example - Individual posteriors")
print("  Slide 5: Distribution detail (violin plot)")

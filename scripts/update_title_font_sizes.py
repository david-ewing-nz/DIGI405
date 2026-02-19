"""Update title font size in presentation/*.pptx (excluding Final_Presentation and temp)"""

from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Pt

TARGET_SIZE = Pt(24)


def is_title_placeholder(shape):
    try:
        if not shape.is_placeholder:
            return False
        ph_type = shape.placeholder_format.type
        return ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    except Exception:
        return False


def set_title_font_size(slide):
    updated = False
    max_size = None
    max_runs = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        is_title = False
        if is_title_placeholder(shape):
            is_title = True
        elif shape.name and "title" in shape.name.lower():
            is_title = True
        else:
            # Heuristic: text box near top of slide
            try:
                if shape.top <= 914400:  # 1 inch in EMU
                    is_title = True
            except Exception:
                pass

        if is_title:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = TARGET_SIZE
            # If no runs, set default for paragraph
            if shape.text_frame.paragraphs:
                shape.text_frame.paragraphs[0].font.size = TARGET_SIZE
            updated = True
        else:
            # Track largest font size as fallback
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size is None:
                        continue
                    if max_size is None or run.font.size > max_size:
                        max_size = run.font.size
                        max_runs = [run]
                    elif run.font.size == max_size:
                        max_runs.append(run)

    if not updated and max_runs:
        for run in max_runs:
            run.font.size = TARGET_SIZE
        updated = True
    return updated


def update_file(pptx_path: Path):
    prs = Presentation(str(pptx_path))
    changed = False
    for slide in prs.slides:
        if set_title_font_size(slide):
            changed = True
    if changed:
        prs.save(str(pptx_path))
    return changed


def main():
    presentation_dir = Path("presentation")
    excluded = {"Final_Presentation.pptx"}
    pptx_files = [
        p
        for p in presentation_dir.glob("*.pptx")
        if p.name not in excluded and not p.name.startswith("~$")
    ]

    updated = []
    skipped = []
    for pptx in pptx_files:
        if update_file(pptx):
            updated.append(pptx.name)
        else:
            skipped.append(pptx.name)

    print("Updated:")
    for name in sorted(updated):
        print(f"  {name}")
    if skipped:
        print("Skipped (no title found):")
        for name in sorted(skipped):
            print(f"  {name}")


if __name__ == "__main__":
    main()

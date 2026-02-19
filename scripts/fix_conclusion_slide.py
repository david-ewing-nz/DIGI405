"""Fix Conclusion_Slide spacing to fit on page"""

from pptx import Presentation
from pptx.util import Pt

prs = Presentation("presentation/Conclusion_Slide.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        text_frame = shape.text_frame
        # Reduce line spacing to make content fit
        text_frame.word_wrap = True
        
        for paragraph in text_frame.paragraphs:
            # Reduce paragraph spacing
            paragraph.space_before = Pt(6)
            paragraph.space_after = Pt(6)
            
            # Slightly reduce font sizes
            for run in paragraph.runs:
                if run.font.size and run.font.size > Pt(14):
                    if run.font.size == Pt(32):  # Title
                        run.font.size = Pt(28)
                    elif run.font.size == Pt(18):  # Section headers
                        run.font.size = Pt(16)
                    elif run.font.size == Pt(14):  # Body
                        run.font.size = Pt(12)

prs.save("presentation/Conclusion_Slide.pptx")
print("Fixed Conclusion_Slide.pptx spacing")

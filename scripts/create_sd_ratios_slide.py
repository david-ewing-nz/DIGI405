"""
Create SD Ratios heatmap slide as native PowerPoint table
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# Add blank slide
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Add title
left = Inches(0.5)
top = Inches(0.3)
width = Inches(9)
height = Inches(0.5)
title_box = slide.shapes.add_textbox(left, top, width, height)
title_frame = title_box.text_frame
title_frame.text = "SD Ratios: VB / Gibbs"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(28)
title_para.font.bold = True
title_para.alignment = PP_ALIGN.CENTER

# Add subtitle
left = Inches(0.5)
top = Inches(0.75)
width = Inches(9)
height = Inches(0.3)
subtitle_box = slide.shapes.add_textbox(left, top, width, height)
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Values < 1 indicate under-dispersion"
subtitle_para = subtitle_frame.paragraphs[0]
subtitle_para.font.size = Pt(14)
subtitle_para.alignment = PP_ALIGN.CENTER

# SD ratio data
data = [
    ["M1", "—", "1.000", "0.993", "0.988", "0.994", "NA", "NA", "NA"],
    ["M2_Q5", "5", "0.721", "0.994", "1.014", "0.995", "0.817", "NA", "NA"],
    ["M2_Q10", "10", "0.889", "0.998", "0.994", "0.970", "0.850", "NA", "NA"],
    ["M2_Q20", "20", "0.961", "0.993", "0.988", "0.962", "0.801", "NA", "NA"],
    ["M2_Q50", "50", "0.992", "0.998", "0.985", "0.913", "0.658", "NA", "NA"],
    ["M2_Q100", "100", "0.994", "1.000", "0.986", "0.803", "0.372", "NA", "NA"]
]

# Table position and size
left = Inches(1.5)
top = Inches(1.3)
width = Inches(7)
height = Inches(3.5)

# Create table (7 rows x 9 cols: header + 6 data rows)
table = slide.shapes.add_table(7, 9, left, top, width, height).table

# Set column widths
col_widths = [Inches(1.2), Inches(0.4)] + [Inches(0.65)] * 7
for i, w in enumerate(col_widths):
    table.columns[i].width = w

# Header row
headers = ["Model", "Q", "β₀", "β₁", "β₂", "τₑ", "τᵤ", "σ²ₑ", "σ²ᵤ"]
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.text_frame.paragraphs[0].font.size = Pt(12)
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(200, 200, 200)

# Color mapping function
def get_color(value_str):
    if value_str == "NA" or value_str == "—":
        return RGBColor(220, 220, 220)  # Grey
    try:
        value = float(value_str)
        if value >= 0.95:
            return RGBColor(255, 255, 153)  # Yellow
        elif value >= 0.85:
            return RGBColor(255, 200, 100)  # Orange
        elif value >= 0.70:
            return RGBColor(255, 150, 80)   # Red-orange
        else:
            return RGBColor(200, 50, 50)    # Dark red
    except:
        return RGBColor(220, 220, 220)

# Fill data rows
for row_idx, row_data in enumerate(data, start=1):
    for col_idx, value in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = value
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = 1  # Middle
        
        # Apply color to data cells (not Model or Q columns)
        if col_idx >= 2:
            cell.fill.solid()
            cell.fill.fore_color.rgb = get_color(value)

# Save
output_path = "presentation/SD_Ratios_Table_Slide.pptx"
prs.save(output_path)
print(f"✓ SD Ratios table slide created: {output_path}")

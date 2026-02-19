"""Fix VI_Definition_Slide by removing footer and further reducing spacing"""

from pptx import Presentation
from pptx.util import Pt

prs = Presentation("presentation/VI_Definition_Slide.pptx")

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        text_frame = shape.text_frame
        
        # Remove or clear italicized footer text
        for para in text_frame.paragraphs:
            text = para.text.strip()
            if "transforms a difficult" in text.lower():
                # Clear this paragraph
                para.text = ""
        
        # Tighten spacing on all paragraphs
        for para in text_frame.paragraphs:
            para.space_before = Pt(0)
            para.space_after = Pt(2)
        
        # Reduce all font sizes by 2pt
        for para in text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    current = run.font.size.pt
                    run.font.size = Pt(max(10, current - 2))

prs.save("presentation/VI_Definition_Slide.pptx")
print("Fixed VI_Definition_Slide.pptx")
